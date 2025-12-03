#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Beautify the block_lecture.pptx using templates.
This script applies template designs to the existing presentation.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def copy_text_content(source_shape, target_shape):
    """Copy text content from source shape to target shape"""
    if not source_shape.has_text_frame or not target_shape.has_text_frame:
        return
    
    source_tf = source_shape.text_frame
    target_tf = target_shape.text_frame
    
    # Clear target text frame
    target_tf.clear()
    
    # Copy paragraphs
    for source_para in source_tf.paragraphs:
        if target_tf.paragraphs:
            target_para = target_tf.paragraphs[0] if len(target_tf.paragraphs) == 1 else target_tf.add_paragraph()
        else:
            target_para = target_tf.add_paragraph()
        
        target_para.text = source_para.text
        target_para.level = source_para.level
        target_para.alignment = source_para.alignment
        
        # Copy font properties
        if source_para.runs:
            for run in source_para.runs:
                target_run = target_para.runs[0] if target_para.runs else None
                if target_run and run.font.size:
                    target_run.font.size = run.font.size
                if target_run and run.font.bold is not None:
                    target_run.font.bold = run.font.bold
                if target_run and run.font.color.rgb:
                    target_run.font.color.rgb = run.font.color.rgb

def apply_template1(source_pptx, template_pptx, output_pptx):
    """Apply template.pptx design to the source presentation"""
    print(f"Applying template.pptx to {source_pptx}...")
    
    # Load source and template
    source_prs = Presentation(source_pptx)
    template_prs = Presentation(template_pptx)
    
    # Create new presentation using the template
    new_prs = Presentation(template_pptx)
    
    # Remove all slides from template
    while len(new_prs.slides) > 0:
        rId = new_prs.slides._sldIdLst[0].rId
        new_prs.part.drop_rel(rId)
        del new_prs.slides._sldIdLst[0]
    
    # Process each slide from source
    for slide_idx, source_slide in enumerate(source_prs.slides):
        print(f"Processing slide {slide_idx + 1}/{len(source_prs.slides)}...")
        
        # Choose appropriate layout based on slide content
        if slide_idx == 0:
            # Title slide - use 封面-01
            layout = new_prs.slide_layouts[0]
        else:
            # Content slide - use 标题和内容（一般样式）
            layout = new_prs.slide_layouts[9]
        
        # Add new slide with chosen layout
        new_slide = new_prs.slides.add_slide(layout)
        
        # Collect all text from source slide
        title_text = ""
        content_shapes = []
        
        for shape in source_slide.shapes:
            if shape.has_text_frame and shape.text.strip():
                # For title slide (slide 0), treat the main text box as title
                if slide_idx == 0 and not content_shapes:
                    title_text = shape.text
                # For content slides, check if it's positioned at the top
                elif shape.top < Inches(1.8):
                    title_text = shape.text
                else:
                    content_shapes.append(shape)
        
        # Apply title
        if title_text:
            for new_shape in new_slide.shapes:
                if new_shape.is_placeholder and new_shape.placeholder_format.type == 1:  # Title placeholder
                    new_shape.text = title_text
                    for para in new_shape.text_frame.paragraphs:
                        para.alignment = PP_ALIGN.CENTER if slide_idx == 0 else PP_ALIGN.LEFT
                        if para.runs:
                            para.runs[0].font.bold = True
                    break
        
        # Apply content
        if content_shapes:
            for new_shape in new_slide.shapes:
                if new_shape.is_placeholder and new_shape.placeholder_format.type == 2:  # Content placeholder
                    new_shape.text_frame.clear()
                    # Copy all paragraphs from all content shapes
                    for content_shape in content_shapes:
                        for src_para in content_shape.text_frame.paragraphs:
                            para = new_shape.text_frame.add_paragraph()
                            para.text = src_para.text
                            para.level = src_para.level
                            # Preserve font size if available
                            if src_para.runs and src_para.runs[0].font.size:
                                for run in para.runs:
                                    run.font.size = src_para.runs[0].font.size
                    # Remove the first empty paragraph if it exists
                    if len(new_shape.text_frame.paragraphs) > 0 and new_shape.text_frame.paragraphs[0].text == "":
                        p = new_shape.text_frame.paragraphs[0]._element
                        p.getparent().remove(p)
                    break
    
    # Save the new presentation
    new_prs.save(output_pptx)
    print(f"Saved beautified presentation to: {output_pptx}")
    print(f"Total slides: {len(new_prs.slides)}")

def apply_template2(source_pptx, template_pptx, output_pptx):
    """Apply Template2.pptx design to the source presentation"""
    print(f"Applying Template2.pptx to {source_pptx}...")
    
    # Load source and template
    source_prs = Presentation(source_pptx)
    template_prs = Presentation(template_pptx)
    
    # Create new presentation using the template
    new_prs = Presentation(template_pptx)
    
    # Remove all slides from template
    while len(new_prs.slides) > 0:
        rId = new_prs.slides._sldIdLst[0].rId
        new_prs.part.drop_rel(rId)
        del new_prs.slides._sldIdLst[0]
    
    # Process each slide from source
    for slide_idx, source_slide in enumerate(source_prs.slides):
        print(f"Processing slide {slide_idx + 1}/{len(source_prs.slides)}...")
        
        # Choose appropriate layout based on slide content
        if slide_idx == 0:
            # Title slide - use TITLE layout
            layout = new_prs.slide_layouts[0]
        else:
            # Content slide - use TITLE_AND_BODY layout
            layout = new_prs.slide_layouts[2]
        
        # Add new slide with chosen layout
        new_slide = new_prs.slides.add_slide(layout)
        
        # Collect all text from source slide
        title_text = ""
        content_shapes = []
        
        for shape in source_slide.shapes:
            if shape.has_text_frame and shape.text.strip():
                # For title slide (slide 0), treat the main text box as title
                if slide_idx == 0 and not content_shapes:
                    title_text = shape.text
                # For content slides, check if it's positioned at the top
                elif shape.top < Inches(1.8):
                    title_text = shape.text
                else:
                    content_shapes.append(shape)
        
        # Apply title
        if title_text:
            for new_shape in new_slide.shapes:
                # Handle both CENTER_TITLE (3) and TITLE (1) placeholder types
                if new_shape.is_placeholder and new_shape.placeholder_format.type in [1, 3]:
                    new_shape.text = title_text
                    for para in new_shape.text_frame.paragraphs:
                        para.alignment = PP_ALIGN.CENTER if slide_idx == 0 else PP_ALIGN.LEFT
                        if para.runs:
                            para.runs[0].font.bold = True
                    break
        
        # Apply content
        if content_shapes:
            for new_shape in new_slide.shapes:
                if new_shape.is_placeholder and new_shape.placeholder_format.type == 2:  # Content placeholder
                    new_shape.text_frame.clear()
                    # Copy all paragraphs from all content shapes
                    for content_shape in content_shapes:
                        for src_para in content_shape.text_frame.paragraphs:
                            para = new_shape.text_frame.add_paragraph()
                            para.text = src_para.text
                            para.level = src_para.level
                            # Preserve font size if available
                            if src_para.runs and src_para.runs[0].font.size:
                                for run in para.runs:
                                    run.font.size = src_para.runs[0].font.size
                    # Remove the first empty paragraph if it exists
                    if len(new_shape.text_frame.paragraphs) > 0 and new_shape.text_frame.paragraphs[0].text == "":
                        p = new_shape.text_frame.paragraphs[0]._element
                        p.getparent().remove(p)
                    break
    
    # Save the new presentation
    new_prs.save(output_pptx)
    print(f"Saved beautified presentation to: {output_pptx}")
    print(f"Total slides: {len(new_prs.slides)}")

def main():
    """Main function to generate beautified presentations"""
    script_dir = os.path.dirname(os.path.abspath(__file__))
    
    source_pptx = os.path.join(script_dir, 'block_lecture.pptx')
    template1_pptx = os.path.join(script_dir, 'template.pptx')
    template2_pptx = os.path.join(script_dir, 'Template2.pptx')
    
    output1_pptx = os.path.join(script_dir, 'block_lecture_beautified_v1.pptx')
    output2_pptx = os.path.join(script_dir, 'block_lecture_beautified_v2.pptx')
    
    # Apply both templates
    print("=" * 60)
    print("Creating beautified version 1 using template.pptx...")
    print("=" * 60)
    apply_template1(source_pptx, template1_pptx, output1_pptx)
    
    print("\n" + "=" * 60)
    print("Creating beautified version 2 using Template2.pptx...")
    print("=" * 60)
    apply_template2(source_pptx, template2_pptx, output2_pptx)
    
    print("\n" + "=" * 60)
    print("✅ Beautification complete!")
    print(f"✅ Created: {output1_pptx}")
    print(f"✅ Created: {output2_pptx}")
    print("=" * 60)

if __name__ == "__main__":
    main()
