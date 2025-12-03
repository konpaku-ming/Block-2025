#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Convert mathematical content in block_lecture_beautified_v2.pptx to LaTeX-style formulas.
This script will enhance mathematical expressions by converting Unicode math symbols to 
proper LaTeX notation within the PowerPoint text.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import re
import os

def convert_math_to_latex(text):
    """
    Convert Unicode mathematical notation to LaTeX-style notation.
    This function wraps mathematical expressions in $...$ and converts symbols.
    """
    # Don't modify empty text
    if not text or not text.strip():
        return text
    
    result = text
    
    # Step 1: Convert Unicode math symbols to LaTeX commands (without wrapping in $)
    # Subscripts - convert to LaTeX
    result = re.sub(r'([a-zA-Z])₀', r'\1_0', result)
    result = re.sub(r'([a-zA-Z])₁', r'\1_1', result)
    result = re.sub(r'([a-zA-Z])₂', r'\1_2', result)
    result = re.sub(r'([a-zA-Z])₃', r'\1_3', result)
    result = re.sub(r'([a-zA-Z])₄', r'\1_4', result)
    result = re.sub(r'([a-zA-Z])₅', r'\1_5', result)
    result = re.sub(r'([a-zA-Z])₆', r'\1_6', result)
    result = re.sub(r'([a-zA-Z])₇', r'\1_7', result)
    result = re.sub(r'([a-zA-Z])₈', r'\1_8', result)
    result = re.sub(r'([a-zA-Z])₉', r'\1_9', result)
    result = re.sub(r'([a-zA-Z])ₙ', r'\1_n', result)
    result = re.sub(r'([a-zA-Z])ᵢ', r'\1_i', result)
    result = re.sub(r'([a-zA-Z])ᵣ', r'\1_r', result)
    
    # Superscripts
    result = re.sub(r'([0-9]+)⁹', r'\1^9', result)
    
    # Square root - replace with LaTeX
    result = re.sub(r'√([a-zA-Z0-9_]+)', r'\\sqrt{\1}', result)
    result = re.sub(r'√', r'\\sqrt{n}', result)
    
    # Floor brackets
    result = re.sub(r'⌊', r'\\lfloor ', result)
    result = re.sub(r'⌋', r' \\rfloor', result)
    
    # Comparison operators
    result = re.sub(r'≤', r' \\leq ', result)
    result = re.sub(r'≥', r' \\geq ', result)
    result = re.sub(r'≠', r' \\neq ', result)
    
    # Arrows
    result = re.sub(r'←', r' \\leftarrow ', result)
    result = re.sub(r'→', r' \\rightarrow ', result)
    
    # Step 2: Protect complex expressions by replacing them with placeholders
    # This prevents nested wrapping
    
    # Protect O() expressions
    o_expressions = []
    def save_o_expr(match):
        o_expressions.append(match.group(1))
        return f'@@OBIG{len(o_expressions)-1}@@'
    result = re.sub(r'O\(([^)]+)\)', save_o_expr, result)
    
    # Protect floor expressions
    floor_expressions = []
    def save_floor_expr(match):
        floor_expressions.append(match.group(1))
        return f'@@FLOOR{len(floor_expressions)-1}@@'
    result = re.sub(r'\\lfloor\s+([^\\]+?)\s+\\rfloor', save_floor_expr, result)
    
    # Step 3: Now wrap mathematical expressions in $ ... $ (won't affect protected contents)
    
    # Wrap sqrt expressions
    result = re.sub(r'\\sqrt\{([a-zA-Z0-9_]+)\}', r'$\\sqrt{\1}$', result)
    
    # Wrap variables with subscripts (e.g., a_1, a_i, a_n)
    result = re.sub(r'\b([a-zA-Z])_\{([a-zA-Z0-9]+)\}', r'$\1_{\2}$', result)
    result = re.sub(r'\b([a-zA-Z])_([a-zA-Z0-9]+)\b', r'$\1_{\2}$', result)
    
    # Wrap superscripts  
    result = re.sub(r'\b([0-9]+)\^([0-9]+)\b', r'$\1^{\2}$', result)
    
    # Wrap comparison operators (with surrounding context)
    result = re.sub(r'(\s)\\(leq|geq|neq)(\s)', r'$\\\2$', result)
    
    # Wrap arrows (with surrounding context)
    result = re.sub(r'(\s)\\(leftarrow|rightarrow)(\s)', r'$\\\2$', result)
    
    # Step 4: Restore protected expressions with wrapping
    for i, expr in enumerate(o_expressions):
        result = result.replace(f'@@OBIG{i}@@', f'$O({expr})$')
    
    for i, expr in enumerate(floor_expressions):
        result = result.replace(f'@@FLOOR{i}@@', f'$\\lfloor {expr} \\rfloor$')
    
    # Step 5: Clean up - merge adjacent math mode markers
    # Remove empty math mode
    result = re.sub(r'\$\s*\$', ' ', result)
    # Merge consecutive math expressions (run multiple times for complex cases)
    for _ in range(20):
        old_result = result
        result = re.sub(r'\$([^$\n]*?)\$\s*\$([^$\n]*?)\$', r'$\1 \2$', result)
        if old_result == result:
            break
    
    # Add space after commas before $
    result = re.sub(r',\$', ', $', result)
    
    # Clean up extra spaces
    result = re.sub(r'\s+', ' ', result)
    
    return result

def process_presentation(input_file, output_file):
    """
    Process the presentation and convert mathematical content to LaTeX notation.
    """
    print(f"Loading presentation: {input_file}")
    prs = Presentation(input_file)
    
    print(f"Total slides: {len(prs.slides)}")
    
    modified_count = 0
    
    for slide_idx, slide in enumerate(prs.slides):
        print(f"Processing slide {slide_idx + 1}...")
        
        for shape_idx, shape in enumerate(slide.shapes):
            if shape.has_text_frame:
                # Process each paragraph
                for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                    original_text = paragraph.text
                    
                    if not original_text or not original_text.strip():
                        continue
                    
                    # Check if text contains mathematical content
                    has_math = any(char in original_text for char in 
                                 ['₀', '₁', '₂', '₃', '₄', '₅', '₆', '₇', '₈', '₉', 
                                  'ₙ', 'ᵢ', 'ᵣ', '√', '⌊', '⌋', '≤', '≥', '≠', '←', '→', '⁹'])
                    has_math = has_math or 'O(' in original_text
                    
                    if has_math:
                        converted_text = convert_math_to_latex(original_text)
                        
                        if converted_text != original_text:
                            print(f"  Slide {slide_idx + 1}, Shape {shape_idx + 1}, Para {para_idx + 1}:")
                            print(f"    Original: {original_text[:80]}...")
                            print(f"    Converted: {converted_text[:80]}...")
                            
                            # Clear existing runs and create new one with converted text
                            # Preserve formatting of the first run if exists
                            if paragraph.runs:
                                first_run = paragraph.runs[0]
                                font_size = first_run.font.size
                                font_bold = first_run.font.bold
                                font_name = first_run.font.name
                                font_color = first_run.font.color.rgb if first_run.font.color.type == 1 else None
                            else:
                                font_size = None
                                font_bold = None
                                font_name = None
                                font_color = None
                            
                            # Clear paragraph
                            paragraph.clear()
                            
                            # Add converted text
                            run = paragraph.add_run()
                            run.text = converted_text
                            
                            # Restore formatting
                            if font_size:
                                run.font.size = font_size
                            if font_bold is not None:
                                run.font.bold = font_bold
                            if font_name:
                                run.font.name = font_name
                            if font_color:
                                run.font.color.rgb = font_color
                            
                            modified_count += 1
    
    print(f"\nTotal modifications: {modified_count}")
    print(f"Saving modified presentation to: {output_file}")
    prs.save(output_file)
    print("Done!")

def main():
    """Main function"""
    script_dir = os.path.dirname(os.path.abspath(__file__))
    
    # Use the repository directory, not the script directory
    repo_dir = '/home/runner/work/Block-2025/Block-2025'
    
    input_file = os.path.join(repo_dir, 'block_lecture_beautified_v2.pptx')
    output_file = os.path.join(repo_dir, 'block_lecture_beautified_v2_latex.pptx')
    
    if not os.path.exists(input_file):
        print(f"Error: Input file not found: {input_file}")
        return 1
    
    process_presentation(input_file, output_file)
    return 0

if __name__ == "__main__":
    exit(main())
