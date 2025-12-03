#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Generate block_lecture.pptx from the LaTeX content.
This script converts the LaTeX lecture notes into a PowerPoint presentation.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_title_slide(prs):
    """Create the title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
    title_frame = title_box.text_frame
    title_frame.text = "分块算法经典例题讲解"
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

def add_section_title(slide, title):
    """Add a section title to a slide"""
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)  # Dark blue

def add_bullet_text(text_frame, text, level=0, font_size=18):
    """Add a bullet point with specified level"""
    p = text_frame.add_paragraph()
    p.text = text
    p.level = level
    p.font.size = Pt(font_size)

def create_example1_slides(prs):
    """Create slides for Example 1: 区间乘法、区间加法与单点查询"""
    
    # Slide 1: Problem statement
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "例题1：区间乘法、区间加法与单点查询")
    
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    p = text_frame.paragraphs[0]
    p.text = "题目描述"
    p.font.size = Pt(24)
    p.font.bold = True
    
    add_bullet_text(text_frame, "给出一个长度为 n 的数列 a₁, a₂, …, aₙ，以及 n 个操作", 0, 16)
    add_bullet_text(text_frame, "", 0, 12)
    
    add_bullet_text(text_frame, "操作类型：", 0, 20)
    add_bullet_text(text_frame, "操作0（区间加法）：opt = 0, l, r, c", 1, 16)
    add_bullet_text(text_frame, "将区间 [l, r] 中的所有数字都加上 c", 2, 14)
    add_bullet_text(text_frame, "", 0, 12)
    
    add_bullet_text(text_frame, "操作1（区间乘法）：opt = 1, l, r, c", 1, 16)
    add_bullet_text(text_frame, "将区间 [l, r] 中的所有数字都乘以 c", 2, 14)
    add_bullet_text(text_frame, "", 0, 12)
    
    add_bullet_text(text_frame, "操作2（单点查询）：opt = 2, r", 1, 16)
    add_bullet_text(text_frame, "查询 aᵣ 的当前值", 2, 14)
    
    # Slide 2: Solution
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "例题1：解法 - 分块 + 懒标记")
    
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    p = text_frame.paragraphs[0]
    p.text = "核心思想"
    p.font.size = Pt(22)
    p.font.bold = True
    
    add_bullet_text(text_frame, "将长度为 n 的数列分成 √n 个块，每块大小约为 √n", 0, 16)
    add_bullet_text(text_frame, "对每个块维护两个懒标记：", 0, 16)
    add_bullet_text(text_frame, "mul[i]：第 i 块的乘法标记（初始为1）", 1, 15)
    add_bullet_text(text_frame, "add[i]：第 i 块的加法标记（初始为0）", 1, 15)
    add_bullet_text(text_frame, "", 0, 12)
    
    add_bullet_text(text_frame, "操作实现：", 0, 20)
    add_bullet_text(text_frame, "区间加法/乘法：", 1, 16)
    add_bullet_text(text_frame, "对于完整覆盖的块：只更新懒标记", 2, 14)
    add_bullet_text(text_frame, "对于部分覆盖的块：暴力修改每个元素", 2, 14)
    add_bullet_text(text_frame, "", 0, 12)
    
    add_bullet_text(text_frame, "单点查询：", 1, 16)
    add_bullet_text(text_frame, "查询 aᵣ 时，返回 aᵣ × mul[block(r)] + add[block(r)]", 2, 14)
    
    # Slide 3: Complexity
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "例题1：复杂度分析")
    
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    p = text_frame.paragraphs[0]
    p.text = "空间复杂度"
    p.font.size = Pt(22)
    p.font.bold = True
    
    add_bullet_text(text_frame, "原数组：O(n)", 0, 16)
    add_bullet_text(text_frame, "块标记：O(√n) 个块，每块2个标记", 0, 16)
    add_bullet_text(text_frame, "总空间：O(n)", 0, 16)
    add_bullet_text(text_frame, "", 0, 12)
    
    p = text_frame.add_paragraph()
    p.text = "时间复杂度"
    p.font.size = Pt(22)
    p.font.bold = True
    
    add_bullet_text(text_frame, "区间修改操作（加法或乘法）：", 0, 16)
    add_bullet_text(text_frame, "散块（两端不完整的块）：暴力修改，O(√n)", 1, 15)
    add_bullet_text(text_frame, "完整块：只修改标记，最多 O(√n) 个块，每个 O(1)", 1, 15)
    add_bullet_text(text_frame, "单次操作：O(√n)", 1, 15)
    add_bullet_text(text_frame, "", 0, 12)
    
    add_bullet_text(text_frame, "单点查询操作：直接通过下标计算所属块号，应用标记，O(1)", 0, 16)
    add_bullet_text(text_frame, "", 0, 12)
    
    add_bullet_text(text_frame, "总复杂度：O(n√n)", 0, 20)

def create_example2_slides(prs):
    """Create slides for Example 2: 区间查询与区间赋值"""
    
    # Slide 1: Problem statement
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "例题2：区间查询与区间赋值")
    
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    p = text_frame.paragraphs[0]
    p.text = "题目描述"
    p.font.size = Pt(24)
    p.font.bold = True
    
    add_bullet_text(text_frame, "给出一个长度为 n 的数列 a₁, a₂, …, aₙ，以及 n 个操作", 0, 16)
    add_bullet_text(text_frame, "", 0, 12)
    
    add_bullet_text(text_frame, "操作内容（每个操作包含两步）：", 0, 20)
    add_bullet_text(text_frame, "每个操作由三个参数 l, r, c 组成：", 0, 16)
    add_bullet_text(text_frame, "1. 查询：统计区间 [l, r] 中有多少个元素等于 c", 1, 16)
    add_bullet_text(text_frame, "2. 修改：将区间 [l, r] 中的所有元素都赋值为 c", 1, 16)
    add_bullet_text(text_frame, "", 0, 12)
    
    add_bullet_text(text_frame, "示例：", 0, 20)
    add_bullet_text(text_frame, "初始数组：[1, 2, 2, 3, 3]", 0, 16)
    add_bullet_text(text_frame, "操作 l=2, r=4, c=2：", 0, 16)
    add_bullet_text(text_frame, "查询结果：区间 [2,4] 是 [2, 2, 3]，有2个元素等于2", 1, 14)
    add_bullet_text(text_frame, "修改后：[1, 2, 2, 2, 3]", 1, 14)
    
    # Slide 2: Solution
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "例题2：解法 - 分块 + 区间赋值标记")
    
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    p = text_frame.paragraphs[0]
    p.text = "核心思想"
    p.font.size = Pt(22)
    p.font.bold = True
    
    add_bullet_text(text_frame, "将数组分成 √n 个块，每块维护：", 0, 16)
    add_bullet_text(text_frame, "tag[i]：整块赋值标记（-1表示无标记）", 1, 15)
    add_bullet_text(text_frame, "每个元素的实际值", 1, 15)
    add_bullet_text(text_frame, "", 0, 12)
    
    add_bullet_text(text_frame, "操作实现：", 0, 20)
    add_bullet_text(text_frame, "散块（两端不完整块）：", 1, 16)
    add_bullet_text(text_frame, "先下传该块的标记（如果有）", 2, 14)
    add_bullet_text(text_frame, "遍历散块元素，统计等于 c 的个数", 2, 14)
    add_bullet_text(text_frame, "将散块元素逐个赋值为 c", 2, 14)
    add_bullet_text(text_frame, "", 0, 12)
    
    add_bullet_text(text_frame, "整块（完全覆盖的块）：", 1, 16)
    add_bullet_text(text_frame, "如果该块有标记 tag[i]：", 2, 14)
    add_bullet_text(text_frame, "若 tag[i] = c，贡献块大小个 c", 3, 13)
    add_bullet_text(text_frame, "否则贡献0个 c", 3, 13)
    add_bullet_text(text_frame, "否则遍历块内所有元素统计", 2, 14)
    add_bullet_text(text_frame, "给整块打上标记 tag[i] = c", 2, 14)
    
    # Slide 3: Complexity
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "例题2：复杂度分析")
    
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    p = text_frame.paragraphs[0]
    p.text = "单次操作时间复杂度"
    p.font.size = Pt(22)
    p.font.bold = True
    
    add_bullet_text(text_frame, "散块处理：O(B)", 0, 16)
    add_bullet_text(text_frame, "整块标记：O(n/B)", 0, 16)
    add_bullet_text(text_frame, "整块查询（最坏）：O(B × n/B) = O(n)", 0, 16)
    add_bullet_text(text_frame, "单次总复杂度：O(n)", 0, 16)
    add_bullet_text(text_frame, "", 0, 12)
    
    p = text_frame.add_paragraph()
    p.text = "均摊分析"
    p.font.size = Pt(22)
    p.font.bold = True
    
    add_bullet_text(text_frame, "关键观察：整块查询只在块没有标记时遍历", 0, 16)
    add_bullet_text(text_frame, "一旦打上标记后，之后的查询都是 O(1)", 0, 16)
    add_bullet_text(text_frame, "每个元素最多被遍历常数次（打标记前）", 0, 16)
    add_bullet_text(text_frame, "散块操作：O(√n)", 0, 16)
    add_bullet_text(text_frame, "整块标记：O(√n)", 0, 16)
    add_bullet_text(text_frame, "", 0, 12)
    
    add_bullet_text(text_frame, "选择 B = √n，单次操作均摊：O(√n)", 0, 16)
    add_bullet_text(text_frame, "总时间复杂度（均摊）：O(n√n)", 0, 20)

def create_example3_slides(prs):
    """Create slides for Example 3: 区间开方与区间求和"""
    
    # Slide 1: Problem statement
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "例题3：区间开方与区间求和")
    
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    p = text_frame.paragraphs[0]
    p.text = "题目描述"
    p.font.size = Pt(24)
    p.font.bold = True
    
    add_bullet_text(text_frame, "给出一个长度为 n 的数列 a₁, a₂, …, aₙ，以及 n 个操作", 0, 16)
    add_bullet_text(text_frame, "", 0, 12)
    
    add_bullet_text(text_frame, "操作类型：", 0, 20)
    add_bullet_text(text_frame, "操作0（区间开方）：opt = 0, l, r", 1, 16)
    add_bullet_text(text_frame, "对区间 [l, r] 中的每个元素 aᵢ 进行开方", 2, 14)
    add_bullet_text(text_frame, "aᵢ ← ⌊√aᵢ⌋", 2, 14)
    add_bullet_text(text_frame, "", 0, 12)
    
    add_bullet_text(text_frame, "操作1（区间求和）：opt = 1, l, r", 1, 16)
    add_bullet_text(text_frame, "查询区间 [l, r] 中所有元素的和", 2, 14)
    add_bullet_text(text_frame, "", 0, 12)
    
    add_bullet_text(text_frame, "示例：", 0, 20)
    add_bullet_text(text_frame, "初始：[16, 9, 4, 1]", 0, 16)
    add_bullet_text(text_frame, "操作0，l=1, r=3：[4, 3, 2, 1]", 0, 16)
    add_bullet_text(text_frame, "操作1，l=1, r=4：输出 4+3+2+1=10", 0, 16)
    
    # Slide 2: Solution
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "例题3：解法 - 分块 + 区间和维护")
    
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    p = text_frame.paragraphs[0]
    p.text = "核心思想"
    p.font.size = Pt(22)
    p.font.bold = True
    
    add_bullet_text(text_frame, "关键性质：开方操作使数字快速减小", 0, 16)
    add_bullet_text(text_frame, "例：10⁹ → 31622 → 177 → 13 → 3 → 1", 1, 14)
    add_bullet_text(text_frame, "约 log log n 次后变为1", 1, 14)
    add_bullet_text(text_frame, "", 0, 12)
    
    add_bullet_text(text_frame, "将数组分成 √n 个块，每块维护：", 0, 16)
    add_bullet_text(text_frame, "sum[i]：第 i 块的元素和", 1, 15)
    add_bullet_text(text_frame, "max[i]：第 i 块的最大值", 1, 15)
    add_bullet_text(text_frame, "", 0, 12)
    
    add_bullet_text(text_frame, "操作实现：", 0, 20)
    add_bullet_text(text_frame, "区间开方：", 1, 16)
    add_bullet_text(text_frame, "对于整块：如果 max[i] ≤ 1，跳过（已收敛）", 2, 14)
    add_bullet_text(text_frame, "否则遍历块内元素，逐个开方，更新 sum[i] 和 max[i]", 2, 14)
    add_bullet_text(text_frame, "散块：直接暴力修改", 2, 14)
    add_bullet_text(text_frame, "", 0, 12)
    
    add_bullet_text(text_frame, "区间求和：", 1, 16)
    add_bullet_text(text_frame, "整块：直接累加 sum[i]，O(1)", 2, 14)
    add_bullet_text(text_frame, "散块：遍历累加，O(√n)", 2, 14)
    
    # Slide 3: Complexity
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "例题3：复杂度分析")
    
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    p = text_frame.paragraphs[0]
    p.text = "空间复杂度"
    p.font.size = Pt(22)
    p.font.bold = True
    
    add_bullet_text(text_frame, "原数组：O(n)", 0, 16)
    add_bullet_text(text_frame, "块信息：O(√n) 个块，每块存储和与最大值", 0, 16)
    add_bullet_text(text_frame, "总空间：O(n)", 0, 16)
    add_bullet_text(text_frame, "", 0, 12)
    
    p = text_frame.add_paragraph()
    p.text = "时间复杂度 - 关键分析"
    p.font.size = Pt(22)
    p.font.bold = True
    
    add_bullet_text(text_frame, "开方操作的收敛性：", 0, 16)
    add_bullet_text(text_frame, "每个元素最多被开方 O(log log V) 次（V为最大值）", 1, 15)
    add_bullet_text(text_frame, "", 0, 12)
    
    add_bullet_text(text_frame, "区间求和：", 0, 16)
    add_bullet_text(text_frame, "整块：O(√n) 个块，每块 O(1)", 1, 15)
    add_bullet_text(text_frame, "散块：O(√n)", 1, 15)
    add_bullet_text(text_frame, "单次：O(√n)", 1, 15)
    add_bullet_text(text_frame, "", 0, 12)
    
    add_bullet_text(text_frame, "区间开方（均摊）：", 0, 16)
    add_bullet_text(text_frame, "单次操作：最坏 O(n)，均摊 O(√n log log V)", 1, 15)
    add_bullet_text(text_frame, "", 0, 12)
    
    add_bullet_text(text_frame, "总时间复杂度：O(n√n log log V)", 0, 20)

def create_example4_slides(prs):
    """Create slides for Example 4: 区间生长与区间计数"""
    
    # Slide 1: Problem statement
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "例题4：区间生长与区间计数")
    
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    p = text_frame.paragraphs[0]
    p.text = "题目描述"
    p.font.size = Pt(24)
    p.font.bold = True
    
    add_bullet_text(text_frame, "有 n 株花，每株花有一个初始高度（≤ 1000 的自然数）", 0, 16)
    add_bullet_text(text_frame, "有两个角色执行 q 个操作：", 0, 16)
    add_bullet_text(text_frame, "Lily White：使花儿生长", 1, 15)
    add_bullet_text(text_frame, "Yuka：统计满足条件的花", 1, 15)
    add_bullet_text(text_frame, "", 0, 12)
    
    add_bullet_text(text_frame, "操作类型：", 0, 20)
    add_bullet_text(text_frame, "操作M（生长）：M l r h", 1, 16)
    add_bullet_text(text_frame, "使区间 [l, r] 内所有花的高度增加 h", 2, 14)
    add_bullet_text(text_frame, "", 0, 12)
    
    add_bullet_text(text_frame, "操作A（询问）：A l r k", 1, 16)
    add_bullet_text(text_frame, "查询区间 [l, r] 内有多少花的高度不低于 k", 2, 14)
    add_bullet_text(text_frame, "（统计满足 aᵢ ≥ k 的花的数量）", 2, 14)
    add_bullet_text(text_frame, "", 0, 12)
    
    add_bullet_text(text_frame, "示例：", 0, 20)
    add_bullet_text(text_frame, "初始：[5, 3, 8, 2]", 0, 16)
    add_bullet_text(text_frame, "M 1 3 2：[7, 5, 10, 2]", 0, 16)
    add_bullet_text(text_frame, "A 1 4 6：输出2（a₁=7 ≥ 6，a₃=10 ≥ 6）", 0, 16)
    
    # Slide 2: Solution
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "例题4：解法 - 分块 + 懒标记 + 块内排序")
    
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    p = text_frame.paragraphs[0]
    p.text = "核心思想"
    p.font.size = Pt(22)
    p.font.bold = True
    
    add_bullet_text(text_frame, "将数组分成 √n 个块，每块维护：", 0, 16)
    add_bullet_text(text_frame, "add[i]：第 i 块的加法懒标记", 1, 15)
    add_bullet_text(text_frame, "sorted[i]：第 i 块元素的有序副本", 1, 15)
    add_bullet_text(text_frame, "", 0, 12)
    
    add_bullet_text(text_frame, "操作实现：", 0, 20)
    add_bullet_text(text_frame, "区间加法（M操作）：", 1, 16)
    add_bullet_text(text_frame, "整块：直接增加 add[i] 标记，O(1)", 2, 14)
    add_bullet_text(text_frame, "散块：", 2, 14)
    add_bullet_text(text_frame, "下传标记到块内所有元素", 3, 13)
    add_bullet_text(text_frame, "逐个修改散块元素", 3, 13)
    add_bullet_text(text_frame, "重新排序该块的有序副本，O(√n log √n)", 3, 13)
    add_bullet_text(text_frame, "", 0, 12)
    
    add_bullet_text(text_frame, "区间计数（A操作）：", 1, 16)
    add_bullet_text(text_frame, "整块：在有序副本中二分查找", 2, 14)
    add_bullet_text(text_frame, "查找第一个 ≥ k - add[i] 的位置", 3, 13)
    add_bullet_text(text_frame, "该位置右侧的元素都满足条件，O(log √n)", 3, 13)
    add_bullet_text(text_frame, "散块：遍历元素，逐个判断，O(√n)", 2, 14)
    
    # Slide 3: Complexity
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "例题4：复杂度分析")
    
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    p = text_frame.paragraphs[0]
    p.text = "空间复杂度"
    p.font.size = Pt(22)
    p.font.bold = True
    
    add_bullet_text(text_frame, "原数组：O(n)", 0, 16)
    add_bullet_text(text_frame, "有序副本：每块 O(√n)，共 O(√n) 个块，总计 O(n)", 0, 16)
    add_bullet_text(text_frame, "懒标记：O(√n)", 0, 16)
    add_bullet_text(text_frame, "总空间：O(n)", 0, 16)
    add_bullet_text(text_frame, "", 0, 12)
    
    p = text_frame.add_paragraph()
    p.text = "时间复杂度"
    p.font.size = Pt(22)
    p.font.bold = True
    
    add_bullet_text(text_frame, "M操作（区间加法）：", 0, 16)
    add_bullet_text(text_frame, "整块标记：O(√n)", 1, 15)
    add_bullet_text(text_frame, "散块修改+重排：O(√n log √n)", 1, 15)
    add_bullet_text(text_frame, "单次：O(√n log √n)", 1, 15)
    add_bullet_text(text_frame, "", 0, 12)
    
    add_bullet_text(text_frame, "A操作（区间计数）：", 0, 16)
    add_bullet_text(text_frame, "整块二分：O(√n log √n)", 1, 15)
    add_bullet_text(text_frame, "散块遍历：O(√n)", 1, 15)
    add_bullet_text(text_frame, "单次：O(√n log √n)", 1, 15)
    add_bullet_text(text_frame, "", 0, 12)
    
    add_bullet_text(text_frame, "q 个操作，每次 O(√n log √n)", 0, 16)
    add_bullet_text(text_frame, "总时间复杂度：O(q√n log n)", 0, 20)

def main():
    """Main function to generate the PowerPoint presentation"""
    # Create a presentation object
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Create title slide
    create_title_slide(prs)
    
    # Create slides for each example
    create_example1_slides(prs)
    create_example2_slides(prs)
    create_example3_slides(prs)
    create_example4_slides(prs)
    
    # Save the presentation
    output_file = '/home/runner/work/Block-2025/Block-2025/block_lecture.pptx'
    prs.save(output_file)
    print(f"PowerPoint presentation saved to: {output_file}")
    print(f"Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
